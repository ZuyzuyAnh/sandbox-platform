#!/usr/bin/env python3
"""Generate Flezi Sandbox introductory PowerPoint (parts 1–3)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Flezi-Sandbox-Intro.pptx"

# Brand palette
NAVY = RGBColor(0x0F, 0x17, 0x2A)
VIOLET = RGBColor(0x6D, 0x28, 0xD9)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT = RGBColor(0x10, 0xB9, 0x81)


def _set_title(shape, text: str, size: int = 32) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = NAVY


def _add_bullets(slide, left, top, width, height, items, font_size=18, sub=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.level = 0
            p.font.size = Pt(font_size)
            p.font.color.rgb = NAVY
            if len(item) > 1:
                for sub_item in item[1]:
                    sp = tf.add_paragraph()
                    sp.text = sub_item
                    sp.level = 1
                    sp.font.size = Pt(font_size - 2)
                    sp.font.color.rgb = SLATE
        else:
            p.text = item
            p.level = 1 if sub else 0
            p.font.size = Pt(font_size)
            p.font.color.rgb = NAVY if not sub else SLATE
    return box


def _footer_refs(slide, refs: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(9), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sources: " + " | ".join(refs)
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE


def _section_slide(prs, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VIOLET
    bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.7))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(22)
        sp.font.color.rgb = SLATE


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.4), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = "Flezi Sandbox"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
    sp = s.text_frame.paragraphs[0]
    sp.text = "Secure, isolated environments for agentic coding at scale"
    sp.font.size = Pt(22)
    sp.font.color.rgb = LIGHT
    m = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8), Inches(0.5))
    mp = m.text_frame.paragraphs[0]
    mp.text = "Business & technical overview  |  Parts 1–3"
    mp.font.size = Pt(14)
    mp.font.color.rgb = SLATE

    # ═══════════════════════════════════════
    # PART 1 — THE PROBLEM
    # ═══════════════════════════════════════
    _section_slide(prs, "Part 1", "The problem today — and why it matters")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Agentic coding is now production reality")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(5),
        [
            "Developers delegate end-to-end work to AI agents: read repos, edit files, run shells, call APIs, deploy.",
            "Tools like Claude Code, Cursor, Copilot, and custom MCP stacks multiply speed — and blast radius.",
            "The new unit of risk is not a bad line of code — it is an autonomous action chain on your infrastructure.",
            "Security and platform teams must answer: where does agent code actually run, and what can it reach?",
        ],
        font_size=20,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "What organizations try first — and where it breaks")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(4.8),
        [
            ("Policy & guardrails", [
                "Prompt rules, tool allowlists, “don’t read .env”, human approval on dangerous commands",
            ]),
            ("Content safety", [
                "Input/output filters, PII redaction, jailbreak classifiers on the model",
            ]),
            ("Developer discipline", [
                "Least-privilege laptops, secrets in vaults, separate cloud accounts",
            ]),
            "These reduce accidents — they do not create a hard boundary when the agent executes code.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Why guardrails alone do not fix system vulnerability")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(8.6),
        Inches(4.5),
        [
            "Guardrails operate on intent and text — sandboxes operate on physics (kernel, network, filesystem).",
            "Prompt injection turns policy into suggestions: Microsoft documents RCE paths where “prompts become shells” in agent frameworks.",
            "Agents compound risk: research shows wrapping an LLM in tool-use loops can increase successful attacks (~1.6×) as early refusals are overridden in later steps.",
            "Major coding agents have shipped critical flaws despite guardrails — credential leaks, allowlist bypass, unsandboxed MCP (TrustFall-class issues).",
            "Industry review: every common isolation tier has a known failure mode; guardrails are necessary but insufficient for code execution.",
        ],
        font_size=17,
    )
    _footer_refs(
        slide,
        [
            "Microsoft Security Blog (2026)",
            "arXiv:2510.01359",
            "Adversa AI TrustFall",
            "Pillar Security (2026)",
            "airuntimesecurity.io",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "You still need a sandbox — the convincing case")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(8.6),
        Inches(4.6),
        [
            "Containment: compromised or tricked agent code must not read production secrets, lateral-move, or persist on the host.",
            "Blast-radius cap: one session = one disposable environment; delete the box, delete the incident.",
            "Evidence & audit: per-session boundaries make logging, TTL, and fleet metrics tractable for compliance.",
            "Separation of duties: builders keep their laptops; agents run in an untrusted tier designed for arbitrary code.",
            "Without sandboxing, “secure agentic coding” is policy on a machine that already trusts the agent as root user.",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Symptoms teams feel today")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(4.8),
        [
            "Ad-hoc Docker on developer laptops — inconsistent, ungoverned, hard to observe centrally",
            "Shared CI runners or cloud VMs for experiments — cross-tenant leakage and credential sprawl",
            "No single pane for active agent sessions, CPU, or session lifetime",
            "Security blocks AI adoption entirely — or accepts unbounded risk",
            "→ Need a productized sandbox fleet, not another policy PDF",
        ],
        font_size=19,
    )

    # ═══════════════════════════════════════
    # PART 2 — THE SOLVER
    # ═══════════════════════════════════════
    _section_slide(prs, "Part 2", "The solver — Flezi Sandbox")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "What is Flezi Sandbox?")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(5),
        [
            "A managed sandbox control plane + dashboard for interactive AI development.",
            "Spins up isolated VS Code (code-server) environments per session — in the browser, on your infrastructure.",
            "Built on OpenSandbox (container lifecycle) with a FastAPI backend and Next.js fleet overview.",
            "Designed for teams who want agentic coding speed with infrastructure-grade isolation and observability.",
        ],
        font_size=20,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Concept — how it fits your world")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(4.2),
        Inches(5),
        [
            "User opens Flezi dashboard",
            "Clicks to start a session",
            "Platform provisions container",
            "code-server + workspace ready",
            "User works with AI extensions (e.g. Claude) inside the sandbox",
            "Session TTL expires → environment destroyed",
        ],
        font_size=17,
    )
    # Simple architecture text block
    arch = slide.shapes.add_textbox(Inches(5.1), Inches(1.55), Inches(4.3), Inches(5))
    tf = arch.text_frame
    tf.word_wrap = True
    lines = [
        "Architecture (logical)",
        "",
        "Browser → nginx",
        "    ├─ Dashboard (Next.js)",
        "    └─ API (FastAPI)",
        "           ├─ PostgreSQL / Redis",
        "           └─ OpenSandbox server",
        "                  └─ Docker sandboxes",
        "                       └─ VS Code :8443",
        "",
        "Session URLs use dedicated",
        "host ports — not your laptop.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14 if i > 0 else 16)
        p.font.bold = i == 0
        p.font.color.rgb = NAVY if i == 0 else SLATE

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Business value proposition")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(5),
        [
            "Enable AI-assisted development without giving agents the keys to production",
            "Standardize sandbox images, TTL, and resource limits across the org",
            "Operational visibility: active sessions, pool utilization, activity log",
            "Deploy on your cloud (e.g. AWS EC2) — data and compute stay in your boundary",
            "Path to hardening: TLS, domain, private registry images, pre-baked extensions",
        ],
        font_size=19,
    )

    # ═══════════════════════════════════════
    # PART 3 — TECHNICAL
    # ═══════════════════════════════════════
    _section_slide(prs, "Part 3", "Core technology, features & benefits")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Core technical stack")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(4.3),
        Inches(5.2),
        [
            ("Frontend", ["Next.js dashboard — fleet overview, sessions, metrics"]),
            ("Backend", ["FastAPI — sessions API, pool state, WebSocket events"]),
            ("Orchestration", ["OpenSandbox server — Docker-backed sandbox lifecycle"]),
            ("Runtime", ["code-server in isolated containers; execd for in-sandbox control"]),
            ("Data", ["PostgreSQL sessions; Redis; real-time activity stream"]),
            ("Edge", ["nginx reverse proxy; optional Let’s Encrypt TLS"]),
        ],
        font_size=16,
    )
    _add_bullets(
        slide,
        Inches(5.0),
        Inches(1.55),
        Inches(4.5),
        Inches(5),
        [
            "Session flow (technical)",
            "1. POST /api/sessions → create sandbox",
            "2. Wait for Running → start code-server via execd",
            "3. Return browser URL (host:port/proxy/8443)",
            "4. Poller fetches CPU/memory via execd /metrics",
            "5. DELETE /api/sessions → tear down container",
        ],
        font_size=15,
        sub=True,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Basic features (today)")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(5),
        [
            "One-click VS Code session — full IDE in browser, no local install",
            "Active session list with open-in-new-tab links and terminate control",
            "Fleet dashboard: pool grid, active tasks, CPU bars, activity log",
            "Metrics: active sandboxes, completions, average duration",
            "Configurable sandbox image (VSCODE_IMAGE), CPU/memory limits, session TTL",
            "Production compose stack: nginx + frontend + backend + Postgres + Redis + OpenSandbox",
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "Technical benefits")
    _add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(8.6),
        Inches(4.8),
        [
            ("Isolation", ["Per-session containers; OpenSandbox API not exposed publicly"]),
            ("Defense in depth", ["Guardrails in the IDE + hard boundary at infrastructure"]),
            ("Observability", ["Central poller, WebSocket events, session records in DB"]),
            ("Disposable compute", ["TTL + delete — no long-lived agent footholds"]),
            ("Extensibility", ["Custom images (extensions, CLI tools); same-origin API for simple deploy"]),
        ],
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, "References & further reading")
    refs_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5))
    tf = refs_box.text_frame
    refs = [
        "Pillar Security — “Your AI Agent Will Run Untrusted Code. Now What?” (sandbox tier failure modes)",
        "Microsoft Security Blog — “When prompts become shells: RCE vulnerabilities in AI agent frameworks”",
        "arXiv:2507.06850 — “The Dark Side of LLMs: Agent-based Attacks for Complete Computer Takeover”",
        "arXiv:2510.01359 — Agent jailbreak / increased vulnerability in tool-use loops",
        "Adversa AI — TrustFall: coding-agent MCP / unsandboxed execution issues",
        "Trend Micro — “Unveiling AI Agent Vulnerabilities: Code Execution”",
        "AI Runtime Security — “Sandbox Patterns” (guardrails necessary but insufficient)",
    ]
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {ref}"
        p.font.size = Pt(15)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = VIOLET
    bg.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.5))
    p = t.text_frame.paragraphs[0]
    p.text = "Flezi Sandbox"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = t.text_frame.add_paragraph()
    p2.text = "Agentic coding speed. Infrastructure-grade isolation."
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
